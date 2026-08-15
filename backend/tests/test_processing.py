from __future__ import annotations

from io import BytesIO

from docx import Document as DocxDocument
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter


def _create_course(client: TestClient) -> str:
    response = client.post("/api/v1/courses", json={"name": "Physics I"})
    assert response.status_code == 201
    return response.json()["id"]


def _upload(client: TestClient, course_id: str, filename: str, content: bytes) -> str:
    response = client.post(
        f"/api/v1/courses/{course_id}/documents",
        files={"file": (filename, content, "application/octet-stream")},
    )
    assert response.status_code == 201
    return response.json()["id"]


def test_process_text_document_extracts_chunks_and_classifies(client: TestClient) -> None:
    course_id = _create_course(client)
    document_id = _upload(
        client,
        course_id,
        "lecture-01.txt",
        b"Lecture 1\n\nNewton's laws describe motion.\n\nForce equals mass times acceleration.",
    )

    process_response = client.post(
        f"/api/v1/courses/{course_id}/documents/{document_id}/process"
    )

    assert process_response.status_code == 200
    analysis = process_response.json()
    assert analysis["document_type"] == "lecture"
    assert analysis["unit_count"] == 1
    assert analysis["chunk_count"] == 1
    assert analysis["extracted_characters"] > 20
    assert analysis["needs_ocr"] is False

    content_response = client.get(
        f"/api/v1/courses/{course_id}/documents/{document_id}/content"
    )
    assert content_response.status_code == 200
    content = content_response.json()
    assert content["document"]["status"] == "processed"
    assert content["units"][0]["source_label"] == "document"
    assert "Newton's laws" in content["units"][0]["text"]
    assert content["chunks"][0]["source_label"] == "document"


def test_content_requires_processing(client: TestClient) -> None:
    course_id = _create_course(client)
    document_id = _upload(client, course_id, "notes.md", b"# Notes\nEnergy")

    response = client.get(f"/api/v1/courses/{course_id}/documents/{document_id}/content")

    assert response.status_code == 409


def test_reprocessing_replaces_previous_content(client: TestClient) -> None:
    course_id = _create_course(client)
    document_id = _upload(client, course_id, "notes.txt", b"Study notes about mechanics.")

    first = client.post(f"/api/v1/courses/{course_id}/documents/{document_id}/process")
    second = client.post(f"/api/v1/courses/{course_id}/documents/{document_id}/process")

    assert first.status_code == 200
    assert second.status_code == 200

    content = client.get(
        f"/api/v1/courses/{course_id}/documents/{document_id}/content"
    ).json()
    assert len(content["units"]) == second.json()["unit_count"]
    assert len(content["chunks"]) == second.json()["chunk_count"]


def test_process_powerpoint_preserves_slide_references(client: TestClient) -> None:
    presentation = Presentation()
    for title in ("Kinematics", "Newton's Laws"):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        box.text = title

    stream = BytesIO()
    presentation.save(stream)

    course_id = _create_course(client)
    document_id = _upload(client, course_id, "lecture-slides.pptx", stream.getvalue())

    response = client.post(f"/api/v1/courses/{course_id}/documents/{document_id}/process")
    assert response.status_code == 200
    assert response.json()["unit_count"] == 2

    content = client.get(
        f"/api/v1/courses/{course_id}/documents/{document_id}/content"
    ).json()
    assert [unit["source_label"] for unit in content["units"]] == ["slide 1", "slide 2"]
    assert content["units"][1]["text"] == "Newton's Laws"


def test_process_docx_extracts_paragraphs_and_tables(client: TestClient) -> None:
    document = DocxDocument()
    document.add_heading("Thermodynamics", level=1)
    document.add_paragraph("The first law relates heat, work, and internal energy.")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Symbol"
    table.cell(0, 1).text = "Meaning"

    stream = BytesIO()
    document.save(stream)

    course_id = _create_course(client)
    document_id = _upload(client, course_id, "lecture-notes.docx", stream.getvalue())

    response = client.post(f"/api/v1/courses/{course_id}/documents/{document_id}/process")
    assert response.status_code == 200

    content = client.get(
        f"/api/v1/courses/{course_id}/documents/{document_id}/content"
    ).json()
    text = content["units"][0]["text"]
    assert "Thermodynamics" in text
    assert "first law" in text
    assert "Symbol | Meaning" in text


def test_blank_pdf_is_flagged_for_ocr(client: TestClient) -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    stream = BytesIO()
    writer.write(stream)

    course_id = _create_course(client)
    document_id = _upload(client, course_id, "scanned-exam.pdf", stream.getvalue())

    response = client.post(f"/api/v1/courses/{course_id}/documents/{document_id}/process")

    assert response.status_code == 200
    analysis = response.json()
    assert analysis["unit_count"] == 1
    assert analysis["chunk_count"] == 0
    assert analysis["needs_ocr"] is True

    content = client.get(
        f"/api/v1/courses/{course_id}/documents/{document_id}/content"
    ).json()
    assert content["units"][0]["source_label"] == "page 1"


def test_corrupt_supported_document_marks_processing_failed(client: TestClient) -> None:
    course_id = _create_course(client)
    document_id = _upload(client, course_id, "broken.pdf", b"not actually a pdf")

    response = client.post(f"/api/v1/courses/{course_id}/documents/{document_id}/process")

    assert response.status_code == 422
    document = client.get(
        f"/api/v1/courses/{course_id}/documents/{document_id}"
    ).json()
    assert document["status"] == "failed"
