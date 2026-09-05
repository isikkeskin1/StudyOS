from __future__ import annotations

from io import BytesIO
from types import SimpleNamespace

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches

from app.schemas.tutor import TutorCitationRead
from app.services.tutor_provider import (
    OpenAIResponsesProvider,
    TutorDraft,
    validate_grounded_draft,
)


def _create_course(client: TestClient, name: str = "Physics I") -> str:
    response = client.post("/api/v1/courses", json={"name": name})
    assert response.status_code == 201
    return response.json()["id"]


def _upload_and_process(
    client: TestClient,
    course_id: str,
    filename: str,
    content: bytes,
) -> str:
    upload = client.post(
        f"/api/v1/courses/{course_id}/documents",
        files={"file": (filename, content, "application/octet-stream")},
    )
    assert upload.status_code == 201
    document_id = upload.json()["id"]
    processed = client.post(
        f"/api/v1/courses/{course_id}/documents/{document_id}/process"
    )
    assert processed.status_code == 200
    return document_id


def _citation() -> TutorCitationRead:
    return TutorCitationRead(
        rank=1,
        document_id="doc",
        document_name="lecture.pdf",
        document_type="lecture",
        chunk_id="chunk",
        source_label="page 2",
        locator_type="page",
        locator_index=2,
        source_reference="lecture.pdf — page 2",
        excerpt="Newton's second law states that net force equals mass times acceleration.",
        relevance_score=0.9,
        lexical_score=0.9,
        topic_affinity=0.0,
        term_coverage=1.0,
        matched_terms=["force", "acceleration"],
    )


def test_tutor_search_ranks_relevant_course_chunks(client: TestClient) -> None:
    course_id = _create_course(client)
    _upload_and_process(
        client,
        course_id,
        "mechanics-lecture.txt",
        (
            b"Newton's second law states that net force equals mass times acceleration. "
            b"Acceleration points in the same direction as the net force."
        ),
    )
    _upload_and_process(
        client,
        course_id,
        "thermal-notes.txt",
        b"Temperature measures thermal state. Heat can transfer between systems.",
    )

    response = client.post(
        f"/api/v1/courses/{course_id}/tutor/search",
        json={"query": "net force acceleration", "limit": 4},
    )

    assert response.status_code == 200
    body = response.json()
    assert body["retrieval_model"] == "lexical-bm25-v1"
    assert body["retrieval_components"] == ["bm25"]
    assert body["topic_signal_applied"] is False
    assert body["result_count"] >= 1
    first = body["citations"][0]
    assert first["document_name"] == "mechanics-lecture.txt"
    assert first["source_reference"] == "mechanics-lecture.txt — document"
    assert first["term_coverage"] == 1.0
    assert set(first["matched_terms"]) == {"net", "force", "acceleration"}


def test_tutor_ask_returns_validated_cited_answer_when_supported(client: TestClient) -> None:
    course_id = _create_course(client)
    _upload_and_process(
        client,
        course_id,
        "newton-notes.txt",
        (
            b"Newton's second law states that net force equals mass times acceleration. "
            b"For constant mass, greater net force produces greater acceleration."
        ),
    )

    response = client.post(
        f"/api/v1/courses/{course_id}/tutor/ask",
        json={"question": "What does Newton's second law say about force and acceleration?"},
    )

    assert response.status_code == 200
    body = response.json()
    assert body["grounding_status"] == "supported"
    assert body["answer_mode"] == "grounded-synthesis-v2"
    assert body["provider_requested"] == "auto"
    assert body["synthesis_provider"] == "local-grounded-v1"
    assert body["validation_status"] == "passed"
    assert body["validation_model"] == "citation-overlap-v2"
    assert body["citation_coverage"] == 1.0
    assert body["grounding_score"] >= body["minimum_claim_support"]
    assert body["validated_claim_count"] >= 1
    assert body["unsupported_claim_count"] == 0
    assert "force" in body["answer"].lower()
    assert "[1]" in body["answer"]
    assert len(body["citations"]) >= 1
    assert body["citations"][0]["source_reference"] == "newton-notes.txt — document"


def test_tutor_refuses_question_without_course_support(client: TestClient) -> None:
    course_id = _create_course(client)
    _upload_and_process(
        client,
        course_id,
        "mechanics.txt",
        b"Velocity is the rate of change of position with respect to time.",
    )

    response = client.post(
        f"/api/v1/courses/{course_id}/tutor/ask",
        json={"question": "Explain quantum entanglement and Bell inequality."},
    )

    assert response.status_code == 200
    body = response.json()
    assert body["grounding_status"] == "insufficient_evidence"
    assert body["validation_status"] == "not_run"
    assert body["synthesis_provider"] == "not_run"
    assert body["citation_coverage"] == 0.0
    assert body["citations"] == []
    assert "couldn't find enough support" in body["answer"].lower()


def test_tutor_preserves_slide_citation_and_course_isolation(client: TestClient) -> None:
    physics_id = _create_course(client, "Physics I")
    other_id = _create_course(client, "Other Course")

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(2))
    box.text = "Impulse equals the change in momentum and is the integral of force over time."
    stream = BytesIO()
    presentation.save(stream)
    _upload_and_process(client, physics_id, "momentum-slides.pptx", stream.getvalue())

    _upload_and_process(
        client,
        other_id,
        "private-other-course.txt",
        b"Impulse is discussed here but must not leak into the Physics I search.",
    )

    physics = client.post(
        f"/api/v1/courses/{physics_id}/tutor/search",
        json={"query": "impulse momentum force"},
    )
    assert physics.status_code == 200
    citations = physics.json()["citations"]
    assert citations[0]["document_name"] == "momentum-slides.pptx"
    assert citations[0]["source_label"] == "slide 1"
    assert citations[0]["locator_type"] == "slide"
    assert citations[0]["locator_index"] == 1
    assert citations[0]["source_reference"] == "momentum-slides.pptx — slide 1"
    assert all(item["document_name"] != "private-other-course.txt" for item in citations)


def test_tutor_hybrid_retrieval_uses_course_topic_evidence(client: TestClient) -> None:
    course_id = _create_course(client)
    _upload_and_process(
        client,
        course_id,
        "momentum.txt",
        (
            b"Momentum Conservation\n"
            b"Momentum conservation states that total momentum remains constant in an isolated "
            b"system. Momentum conservation is useful for collision problems."
        ),
    )
    analysis = client.post(f"/api/v1/courses/{course_id}/analyze")
    assert analysis.status_code == 200

    response = client.post(
        f"/api/v1/courses/{course_id}/tutor/search",
        json={"query": "momentum conservation"},
    )
    assert response.status_code == 200
    body = response.json()
    assert body["retrieval_model"] == "hybrid-topic-bm25-v1"
    assert body["retrieval_components"] == ["bm25", "course_topic_evidence"]
    assert body["topic_signal_applied"] is True
    assert body["citations"][0]["topic_affinity"] > 0
    assert body["citations"][0]["matched_topics"]


def test_grounding_validator_rejects_cited_but_unsupported_claim() -> None:
    draft = TutorDraft(
        answer=(
            "Net force equals mass times acceleration [1]. "
            "Quantum tunneling determines the acceleration direction [1]."
        ),
        provider="test-provider",
    )

    validation = validate_grounded_draft(draft, [_citation()])

    assert validation.status == "rejected"
    assert validation.model == "citation-overlap-v2"
    assert validation.validated_claim_count == 1
    assert validation.unsupported_claim_count == 1
    assert validation.citation_coverage == 0.5
    assert validation.grounding_score < 1.0


def test_openai_provider_uses_grounded_untrusted_source_prompt() -> None:
    class FakeResponses:
        def __init__(self) -> None:
            self.kwargs: dict[str, object] = {}

        def create(self, **kwargs: object) -> SimpleNamespace:
            self.kwargs = kwargs
            return SimpleNamespace(
                output_text="Net force equals mass times acceleration [1]."
            )

    fake_responses = FakeResponses()
    fake_client = SimpleNamespace(responses=fake_responses)
    provider = OpenAIResponsesProvider(
        api_key="test-key",
        model="gpt-test",
        client=fake_client,
    )

    draft = provider.synthesize(
        "What is Newton's second law?",
        [_citation()],
        "guided",
    )

    assert draft.provider == "openai-responses:gpt-test"
    assert draft.insufficient_evidence is False
    assert "[1]" in draft.answer
    instructions = str(fake_responses.kwargs["instructions"])
    assert "untrusted data" in instructions
    assert "only from the supplied course-source packet" in instructions
    assert fake_responses.kwargs["store"] is False


def test_explicit_openai_provider_requires_api_key(client: TestClient) -> None:
    course_id = _create_course(client)
    _upload_and_process(
        client,
        course_id,
        "newton.txt",
        b"Newton's second law states that net force equals mass times acceleration.",
    )

    response = client.post(
        f"/api/v1/courses/{course_id}/tutor/ask",
        json={
            "question": "What does Newton's second law say about force?",
            "provider": "openai",
        },
    )

    assert response.status_code == 503
    assert "OPENAI_API_KEY" in response.json()["detail"]
