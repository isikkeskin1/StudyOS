from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from docx import Document as DocxDocument
from pypdf import PdfReader
from pptx import Presentation


class DocumentExtractionError(RuntimeError):
    pass


@dataclass(frozen=True)
class ExtractedUnit:
    locator_type: str
    locator_index: int | None
    source_label: str
    text: str


def _clean_text(text: str) -> str:
    lines = [line.rstrip() for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")]
    cleaned: list[str] = []
    blank = False
    for line in lines:
        stripped = line.strip()
        if not stripped:
            if cleaned and not blank:
                cleaned.append("")
            blank = True
            continue
        cleaned.append(stripped)
        blank = False
    return "\n".join(cleaned).strip()


def _extract_plain_text(path: Path) -> list[ExtractedUnit]:
    text = path.read_text(encoding="utf-8-sig", errors="replace")
    return [
        ExtractedUnit(
            locator_type="document",
            locator_index=None,
            source_label="document",
            text=_clean_text(text),
        )
    ]


def _extract_pdf(path: Path) -> list[ExtractedUnit]:
    reader = PdfReader(str(path))
    units: list[ExtractedUnit] = []
    for page_number, page in enumerate(reader.pages, start=1):
        units.append(
            ExtractedUnit(
                locator_type="page",
                locator_index=page_number,
                source_label=f"page {page_number}",
                text=_clean_text(page.extract_text() or ""),
            )
        )
    return units


def _extract_docx(path: Path) -> list[ExtractedUnit]:
    document = DocxDocument(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))

    return [
        ExtractedUnit(
            locator_type="document",
            locator_index=None,
            source_label="document",
            text=_clean_text("\n".join(parts)),
        )
    ]


def _extract_pptx(path: Path) -> list[ExtractedUnit]:
    presentation = Presentation(str(path))
    units: list[ExtractedUnit] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = getattr(shape, "text", "")
                if text and text.strip():
                    parts.append(text)

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))

        units.append(
            ExtractedUnit(
                locator_type="slide",
                locator_index=slide_number,
                source_label=f"slide {slide_number}",
                text=_clean_text("\n".join(parts)),
            )
        )

    return units


def extract_document(path: Path, extension: str) -> list[ExtractedUnit]:
    extension = extension.lower()

    try:
        if extension in {".txt", ".md"}:
            return _extract_plain_text(path)
        if extension == ".pdf":
            return _extract_pdf(path)
        if extension == ".docx":
            return _extract_docx(path)
        if extension == ".pptx":
            return _extract_pptx(path)
    except Exception as exc:
        raise DocumentExtractionError(f"Could not extract {extension} document") from exc

    raise DocumentExtractionError(f"No extractor is available for {extension}")
